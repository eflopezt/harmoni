"""
Generar PPTX de presentación Harmoni — Mayo 2026.

Estructura derivada de docs/internal/GUIA_PRESENTACION_DESIGN.md.

Uso:
    python docs/generar_ppt_demo.py
    → genera presentacion/Harmoni_Demo_Mayo_2026.pptx
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ─────────────────────────────────────────────────────────────────────
# Identidad visual Harmoni (de GUIA_PRESENTACION_DESIGN.md §0)
# ─────────────────────────────────────────────────────────────────────

# Colores
TEAL = RGBColor(0x0F, 0x76, 0x6E)      # primario
TEAL_DARK = RGBColor(0x13, 0x4E, 0x4A)
TEAL_LIGHT = RGBColor(0x5E, 0xEA, 0xD4)
SLATE_900 = RGBColor(0x0D, 0x2B, 0x27)
SLATE_700 = RGBColor(0x33, 0x40, 0x49)
SLATE_500 = RGBColor(0x47, 0x55, 0x69)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
WARNING = RGBColor(0xEA, 0xB3, 0x08)
DANGER = RGBColor(0xDC, 0x26, 0x26)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)

# Fuentes
FONT_HEAD = 'Calibri'      # Disponible en todos los Office
FONT_BODY = 'Calibri'
FONT_MONO = 'Consolas'


# ─────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    """Pinta el fondo de la slide del color dado."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=18, bold=False,
             italic=False, color=SLATE_900, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    """Agrega caja de texto en posición/tamaño dado, con estilo."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, color, *, line_color=None, line_width=0):
    """Agrega rectángulo con relleno sólido."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, x, y, w, h, color, *, line_color=None):
    """Agrega un rectángulo redondeado."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_kpi_box(slide, x, y, w, h, value, label, color=TEAL):
    """KPI tile estilo Cockpit V2."""
    add_rounded(slide, x, y, w, h, SLATE_100, line_color=SLATE_300)
    add_text(slide, x, y + Inches(0.15), w, Inches(0.6), value,
             font=FONT_MONO, size=32, bold=True, color=color,
             align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.75), w, Inches(0.3), label,
             font=FONT_BODY, size=10, color=SLATE_500,
             align=PP_ALIGN.CENTER)


def page_footer(slide, num, total):
    """Footer con número de slide y branding."""
    add_text(slide, Inches(0.4), Inches(7.0), Inches(5), Inches(0.25),
             'Harmoni · ERP de RRHH para Perú · Mayo 2026',
             font=FONT_BODY, size=9, color=SLATE_500)
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.8), Inches(0.25),
             f'{num} / {total}',
             font=FONT_BODY, size=9, color=SLATE_500, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────
# Generar presentación
# ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    TOTAL = 18

    # ─────────────────────────────────────────────────────────────────
    # Slide 1 — Portada
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, SLATE_900)

    # Accent corner
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL_LIGHT)

    # Logo (representación textual — el cliente verá el logo real en demo.harmoni.pe)
    add_rounded(s, Inches(5.7), Inches(2.0), Inches(2.0), Inches(2.0), TEAL)
    add_text(s, Inches(5.7), Inches(2.0), Inches(2.0), Inches(2.0), 'H',
             font=FONT_HEAD, size=110, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.8),
             'Harmoni',
             font=FONT_HEAD, size=64, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
             'ERP de Recursos Humanos · Hecho en Perú',
             font=FONT_BODY, size=20, color=TEAL_LIGHT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.4),
             'Presentación demo · Mayo 2026 · v1.2.1',
             font=FONT_BODY, size=14, color=SLATE_300,
             align=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────────
    # Slide 2 — El problema
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), DANGER)

    add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.6),
             'El problema',
             font=FONT_HEAD, size=40, bold=True, color=SLATE_900)

    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
             'RRHH peruano gasta horas en planilla, no en personas.',
             font=FONT_BODY, size=22, color=SLATE_500)

    bullets = [
        ('⏰', '4 horas para cerrar una planilla mensual', 'Excel paralelo, fórmulas frágiles, errores que llegan a SUNAFIL'),
        ('🌎', 'ERPs extranjeros traducidos al español', 'No conocen D.Leg. 728, ni gratificación + 9%, ni CTS semestral'),
        ('📊', 'Datos en silos sin trazabilidad', 'Tareo en Synkro, recibos en PDF, vacaciones en mail, IR 5ta en Excel'),
        ('👤', 'RRHH es el cuello de botella', 'Cada consulta del trabajador pasa por la asistente de RRHH'),
    ]
    y = 2.4
    for icon, title, desc in bullets:
        add_text(s, Inches(0.7), Inches(y), Inches(0.6), Inches(0.5), icon,
                 font=FONT_BODY, size=28, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.4), Inches(y), Inches(11), Inches(0.4), title,
                 font=FONT_HEAD, size=18, bold=True, color=SLATE_900)
        add_text(s, Inches(1.4), Inches(y + 0.4), Inches(11), Inches(0.4), desc,
                 font=FONT_BODY, size=13, color=SLATE_500)
        y += 1.0

    page_footer(s, 2, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 3 — La solución
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.6),
             'La solución: Harmoni',
             font=FONT_HEAD, size=40, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
             'ERP de RRHH peruano · 17 módulos integrados · Diseño Cockpit',
             font=FONT_BODY, size=20, color=SLATE_500)

    pillars = [
        ('Trabaja en armonía', 'Personal → Asistencia → Nóminas → Portal del empleado sin doble ingreso'),
        ('Comunica antes', 'Alertas proactivas, workflows configurables, audit log completo'),
        ('Calcula bien', 'D.Leg. 728/713/650, IR 5ta SUNAT, AFP/ONP, CTS, gratis + 9%'),
    ]
    pw = Inches(3.9)
    px = Inches(0.7)
    py = Inches(2.5)
    gap = Inches(0.2)
    for i, (titulo, desc) in enumerate(pillars):
        x = px + (pw + gap) * i
        add_rounded(s, x, py, pw, Inches(3.5), SLATE_100, line_color=TEAL_LIGHT)
        # Header teal
        add_rounded(s, x, py, pw, Inches(0.7), TEAL)
        add_text(s, x, py, pw, Inches(0.7), titulo,
                 font=FONT_HEAD, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_text(s, x + Inches(0.3), py + Inches(1.0), pw - Inches(0.6), Inches(2.3),
                 desc,
                 font=FONT_BODY, size=15, color=SLATE_700)

    add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
             'En producción hoy: Consorcio Stiler (234 colaboradores) · '
             'Demo: Grupo Sabores del Sur (86 colaboradores · gastronomía)',
             font=FONT_BODY, size=13, color=SLATE_500, align=PP_ALIGN.CENTER)

    page_footer(s, 3, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 4 — Centro de Comando
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, SLATE_100)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Centro de Comando',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Lo que hay que hacer hoy · alertas en vivo · KPIs · "Una cosa para hoy"',
             font=FONT_BODY, size=14, color=SLATE_500)

    # 4 KPIs
    kpi_w = Inches(2.85)
    kpi_h = Inches(1.4)
    kpi_y = Inches(1.7)
    gap = Inches(0.15)
    kpi_x = Inches(0.7)
    add_kpi_box(s, kpi_x + (kpi_w + gap) * 0, kpi_y, kpi_w, kpi_h, '86', 'Colaboradores activos', TEAL)
    add_kpi_box(s, kpi_x + (kpi_w + gap) * 1, kpi_y, kpi_w, kpi_h, '9', 'Locales operando', TEAL)
    add_kpi_box(s, kpi_x + (kpi_w + gap) * 2, kpi_y, kpi_w, kpi_h, 'S/ 220K', 'Planilla est./mes', SUCCESS)
    add_kpi_box(s, kpi_x + (kpi_w + gap) * 3, kpi_y, kpi_w, kpi_h, '10', 'Alertas activas', ACCENT)

    # Alerta top
    alert_y = Inches(3.4)
    add_rounded(s, Inches(0.7), alert_y, Inches(11.9), Inches(1.2), WHITE, line_color=DANGER)
    add_box(s, Inches(0.7), alert_y, Inches(0.15), Inches(1.2), DANGER)
    add_text(s, Inches(1.0), alert_y + Inches(0.15), Inches(11), Inches(0.35),
             '⚡ Una cosa para hacer hoy',
             font=FONT_HEAD, size=11, bold=True, color=DANGER)
    add_text(s, Inches(1.0), alert_y + Inches(0.45), Inches(11), Inches(0.45),
             '10 contratos vencen en próximos 30 días',
             font=FONT_HEAD, size=20, bold=True, color=SLATE_900)
    add_text(s, Inches(1.0), alert_y + Inches(0.85), Inches(11), Inches(0.3),
             'Revisar renovaciones · click para ir al módulo Contratos',
             font=FONT_BODY, size=12, color=SLATE_500)

    # 3 alertas más
    alerts = [
        ('🟡', '8 solicitudes de vacaciones pendientes de aprobación', '/vacaciones/'),
        ('🟠', 'DIGESA vence: 12 trabajadores en 15 días', '/capacitaciones/'),
        ('🟢', '3 nuevos postulantes para Chef Ejecutivo Senior', '/reclutamiento/'),
    ]
    ay = 4.8
    for icon, msg, url in alerts:
        add_rounded(s, Inches(0.7), Inches(ay), Inches(11.9), Inches(0.55), WHITE, line_color=SLATE_300)
        add_text(s, Inches(0.9), Inches(ay + 0.13), Inches(0.5), Inches(0.3), icon, size=14)
        add_text(s, Inches(1.5), Inches(ay + 0.13), Inches(9), Inches(0.3), msg,
                 font=FONT_BODY, size=12, color=SLATE_900)
        add_text(s, Inches(10.5), Inches(ay + 0.13), Inches(2), Inches(0.3), url + ' →',
                 font=FONT_MONO, size=10, color=TEAL, align=PP_ALIGN.RIGHT)
        ay += 0.7

    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             'Frase: "El sistema te dice qué hacer; tú decides cuándo hacerlo."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 4, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 5 — Personal (drawer + bulk + inline)
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Módulo Personal',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Una sola pantalla para gestionar 200+ empleados — sin abrir 5 modales',
             font=FONT_BODY, size=14, color=SLATE_500)

    # Features
    features = [
        ('📋', 'Tabla densa Cockpit', 'Una línea por persona, columnas numéricas en JetBrains Mono'),
        ('🔍', 'Filtros chip + búsqueda', 'Estado, grupo, sede, área, régimen — sin recargar página'),
        ('📑', 'Drawer 540px', 'Click sobre fila abre lateral con 4 tabs (Datos, Contrato, Asistencia, Histórico) sin perder contexto'),
        ('🎯', 'Bulk bar', 'Multi-selección → barra inferior: cambiar estado · asignar área · exportar · cesar masivo'),
        ('✏️', 'Inline edit', 'Doble-click sobre cargo/sueldo → edita en la celda, Enter para confirmar'),
        ('⌘', '⌘K Command Palette', 'Búsqueda global instantánea — empleados, módulos, acciones'),
    ]
    fy = 1.8
    fx = 0.7
    fw = 6.0
    fh = 0.85
    for i, (icon, title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(fx + (fw + 0.3) * col)
        y = Inches(fy + (fh + 0.15) * row)
        add_rounded(s, x, y, Inches(fw), Inches(fh), SLATE_100, line_color=SLATE_300)
        add_text(s, x + Inches(0.15), y + Inches(0.1), Inches(0.5), Inches(0.7), icon,
                 size=24, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.7), y + Inches(0.05), Inches(fw - 0.8), Inches(0.3), title,
                 font=FONT_HEAD, size=13, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.7), y + Inches(0.35), Inches(fw - 0.8), Inches(0.5), desc,
                 font=FONT_BODY, size=10, color=SLATE_500)

    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             'Frase: "Quien usa Harmoni una semana ya no usa el menú — usa ⌘K."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 5, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 6 — Asistencia
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Asistencia',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Matriz mensual · sync biométrico de 1 click · turno rotativo configurable',
             font=FONT_BODY, size=14, color=SLATE_500)

    # Tabla ASCII estilizada
    add_rounded(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.5), SLATE_100, line_color=SLATE_300)
    # Header con días
    header_y = 2.0
    add_text(s, Inches(0.9), Inches(header_y), Inches(2.5), Inches(0.4),
             'Trabajador',
             font=FONT_HEAD, size=11, bold=True, color=SLATE_900)
    dias = ['L', 'M', 'X', 'J', 'V', 'S', 'D', 'L', 'M', 'X', 'J', 'V', 'S', 'D']
    for i, d in enumerate(dias):
        add_text(s, Inches(3.5 + 0.6 * i), Inches(header_y), Inches(0.6), Inches(0.4),
                 d, font=FONT_MONO, size=10, color=SLATE_500, align=PP_ALIGN.CENTER)

    # Filas ejemplo
    estados = [
        ('AGUIRRE CCAMA, E', 'Mozo', ['T','T','DS','T','T','T','T','T','T','DS','T','T','T','T']),
        ('CONDORI APAZA, R', 'Cocinero', ['T','T','T','T','T','DS','T','T','T','T','T','T','DS','T']),
        ('VEGA CORDERO, A', 'Bartender', ['T','T','T','DS','T','T','T','T','T','T','DS','T','T','T']),
        ('FLORES MAMANI, J', 'Mozo Sr', ['T','FA','T','T','T','T','T','VAC','VAC','VAC','VAC','VAC','VAC','VAC']),
        ('TORRES DIAZ, Y', 'Hostess', ['T','T','T','T','DS','T','T','T','T','T','T','DS','T','T']),
    ]
    code_color = {
        'T': SUCCESS,
        'DS': SLATE_500,
        'FA': DANGER,
        'VAC': RGBColor(0x3B, 0x82, 0xF6),  # blue
    }
    for i, (nombre, cargo, codes) in enumerate(estados):
        ry = 2.45 + i * 0.5
        # nombre
        add_text(s, Inches(0.9), Inches(ry), Inches(2.0), Inches(0.25), nombre,
                 font=FONT_HEAD, size=10, bold=True, color=SLATE_900)
        add_text(s, Inches(0.9), Inches(ry + 0.22), Inches(2.0), Inches(0.2), cargo,
                 font=FONT_BODY, size=8, color=SLATE_500)
        # cells
        for j, c in enumerate(codes):
            cx = Inches(3.5 + 0.6 * j)
            cy = Inches(ry)
            cw = Inches(0.55)
            ch = Inches(0.35)
            color = code_color.get(c, SLATE_300)
            add_rounded(s, cx, cy, cw, ch, color)
            add_text(s, cx, cy, cw, ch, c,
                     font=FONT_MONO, size=8, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Leyenda
    ly = 5.6
    leg = [('T', 'Trabajado', SUCCESS), ('DS', 'Descanso semanal (rotativo)', SLATE_500),
           ('FA', 'Falta', DANGER), ('VAC', 'Vacaciones', RGBColor(0x3B, 0x82, 0xF6))]
    for i, (c, lbl, col) in enumerate(leg):
        x = Inches(0.7 + 3.0 * i)
        add_rounded(s, x, Inches(ly), Inches(0.4), Inches(0.3), col)
        add_text(s, x, Inches(ly), Inches(0.4), Inches(0.3), c,
                 font=FONT_MONO, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.45), Inches(ly), Inches(2.5), Inches(0.3), lbl,
                 font=FONT_BODY, size=11, color=SLATE_900, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4),
             '🆕 Turno rotativo: gastronomía no usa roster. El día de descanso semanal ' +
             'es configurable por trabajador (lunes a domingo).',
             font=FONT_BODY, size=12, color=SLATE_700)
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             'Frase: "Sync con tu reloj biométrico de 1 click — sin Excel paralelo."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 6, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 7 — Nóminas — Wizard
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Nóminas · Cierre Guiado',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             '11 pasos · 12 minutos · personalizable por empresa',
             font=FONT_BODY, size=14, color=SLATE_500)

    # 11 steps en grid 2 columnas
    steps = [
        ('1', 'Asistencia del mes', 'fa-clock'),
        ('2', 'Conceptos sin inconsistencias', 'fa-cogs'),
        ('3', 'Crear/abrir período', 'fa-calendar-plus'),
        ('4', 'Generar planilla', 'fa-calculator'),
        ('5', 'Aprobar planilla', 'fa-check-circle'),
        ('6', 'Emitir boletas', 'fa-paper-plane'),
        ('7', 'Acuses de recibo', 'fa-signature'),
        ('8', 'Exportar PLAME a SUNAT', 'fa-file-export'),
        ('9', 'Exportar AFPNet', 'fa-piggy-bank'),
        ('10', 'Asiento contable', 'fa-file-invoice'),
        ('11', 'Cerrar período', 'fa-lock'),
    ]
    sx = 0.7
    sy = 1.8
    sw = 6.0
    sh = 0.4
    for i, (num, title, _) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = Inches(sx + (sw + 0.3) * col)
        y = Inches(sy + (sh + 0.08) * row)
        add_rounded(s, x, y, Inches(sw), Inches(sh), SLATE_100)
        # número
        add_rounded(s, x + Inches(0.1), y + Inches(0.06), Inches(0.3), Inches(0.3), TEAL)
        add_text(s, x + Inches(0.1), y + Inches(0.06), Inches(0.3), Inches(0.3), num,
                 font=FONT_HEAD, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.55), y + Inches(0.05), Inches(sw - 0.65), Inches(0.3), title,
                 font=FONT_BODY, size=12, color=SLATE_900, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
             '🆕 v1.2: AFPNet step #9 nuevo · Botón "Personalizar pasos" (oculta los que no aplican)',
             font=FONT_BODY, size=12, color=SUCCESS, bold=True)
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             'Frase: "Lo que en SAP requiere 4 horas y un Excel paralelo, acá son 12 minutos."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 7, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 8 — Liquidaciones (ESTRELLA)
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, SLATE_100)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Liquidaciones laborales · NUEVO',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Flujo end-to-end al cese · 1 botón · 3 PDFs · 5 etapas workflow',
             font=FONT_BODY, size=14, color=SLATE_500)

    # Flujo horizontal en 5 cajas
    flujo = [
        ('1', 'Iniciar cese', 'Botón en ficha'),
        ('2', 'Wizard 3 pasos', 'Datos · preview · confirmar'),
        ('3', 'Cálculo auto', 'Truncas + descuentos'),
        ('4', 'PDFs', 'Boleta · Carta no adeudo · Certificado'),
        ('5', 'Workflow', 'Offboarding 5 etapas'),
    ]
    fx_start = 0.7
    fw = 2.4
    fh = 1.5
    fy = 1.85
    gap = 0.05
    for i, (num, title, desc) in enumerate(flujo):
        x = Inches(fx_start + (fw + gap) * i)
        y = Inches(fy)
        add_rounded(s, x, y, Inches(fw), Inches(fh), WHITE, line_color=TEAL_LIGHT)
        # número arriba-izq
        add_rounded(s, x + Inches(0.1), y + Inches(0.1), Inches(0.35), Inches(0.35), TEAL)
        add_text(s, x + Inches(0.1), y + Inches(0.1), Inches(0.35), Inches(0.35), num,
                 font=FONT_HEAD, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(fw - 0.4), Inches(0.4), title,
                 font=FONT_HEAD, size=14, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.2), y + Inches(0.95), Inches(fw - 0.4), Inches(0.5), desc,
                 font=FONT_BODY, size=10, color=SLATE_500)
        # arrow next
        if i < 4:
            ax = x + Inches(fw - 0.05)
            add_text(s, ax, y + Inches(fh / 2 - 0.15), Inches(0.4), Inches(0.3), '→',
                     size=24, color=TEAL, align=PP_ALIGN.CENTER)

    # 8 motivos
    add_rounded(s, Inches(0.7), Inches(3.9), Inches(11.9), Inches(1.7), WHITE, line_color=SLATE_300)
    add_text(s, Inches(0.9), Inches(4.05), Inches(11), Inches(0.4),
             '8 motivos soportados:',
             font=FONT_HEAD, size=14, bold=True, color=SLATE_900)
    motivos = ['Renuncia', 'Despido (causa)', 'Despido arbitrario (indemnización)', 'Mutuo disenso',
               'Caducidad contrato', 'Jubilación', 'Fallecimiento', 'Período de prueba']
    mx = 0.9
    my = 4.5
    mw = 2.85
    for i, m in enumerate(motivos):
        col = i % 4
        row = i // 4
        x = Inches(mx + mw * col)
        y = Inches(my + 0.45 * row)
        add_text(s, x, y, Inches(mw), Inches(0.3), '• ' + m,
                 font=FONT_BODY, size=11, color=SLATE_700)

    # KPIs
    kpi_y = 5.85
    add_text(s, Inches(0.7), Inches(kpi_y), Inches(12), Inches(0.4),
             'Calidad: 47 tests passing (modelo + signal + wizard + cartas + offboarding + exit interview)',
             font=FONT_BODY, size=12, color=SUCCESS, bold=True)
    add_text(s, Inches(0.7), Inches(kpi_y + 0.4), Inches(12), Inches(0.3),
             'Frase: "Lo que en SAP requiere 4 horas y un Excel paralelo, acá son 3 clicks."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)
    add_text(s, Inches(0.7), Inches(kpi_y + 0.7), Inches(12), Inches(0.3),
             'URL demo: /personal/<id>/cesar/ · /nominas/liquidacion/<id>/',
             font=FONT_MONO, size=10, color=SLATE_500)

    page_footer(s, 8, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 9 — Agente IA (WOW)
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, SLATE_900)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL_LIGHT)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Agente IA · Harmoni AI',
             font=FONT_HEAD, size=36, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'No es un chatbot. Es un asistente que ACTÚA con tool-use real.',
             font=FONT_BODY, size=16, color=TEAL_LIGHT)

    # Mockup de chat
    chat_x = 0.7
    chat_y = 1.9
    chat_w = 11.9
    chat_h = 4.2
    add_rounded(s, Inches(chat_x), Inches(chat_y), Inches(chat_w), Inches(chat_h), WHITE)

    # Mensaje user
    add_rounded(s, Inches(chat_x + 2), Inches(chat_y + 0.25), Inches(8), Inches(0.55), TEAL)
    add_text(s, Inches(chat_x + 2.2), Inches(chat_y + 0.25), Inches(7.6), Inches(0.55),
             '"El mes pasado olvidé subir el sueldo S/200 a todos los trabajadores"',
             font=FONT_BODY, size=13, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Mensaje IA
    add_rounded(s, Inches(chat_x + 0.3), Inches(chat_y + 1.0), Inches(11.3), Inches(2.6),
                SLATE_100, line_color=TEAL_LIGHT)
    add_text(s, Inches(chat_x + 0.5), Inches(chat_y + 1.1), Inches(0.5), Inches(0.3), '🤖',
             size=20)
    add_text(s, Inches(chat_x + 1.0), Inches(chat_y + 1.1), Inches(10), Inches(0.3),
             'Harmoni AI · usando D.S. 003-97-TR + tool reintegro_masivo',
             font=FONT_HEAD, size=11, bold=True, color=TEAL)
    add_text(s, Inches(chat_x + 0.5), Inches(chat_y + 1.5), Inches(11), Inches(2.0),
             ['Detecté que aplica un reintegro retroactivo para el mes anterior.',
              '',
              '• 64 trabajadores afectados',
              '• Reintegro sueldo: S/ 12,800 (64 × S/200)',
              '• Proporcional gratificación: S/ 2,133',
              '• Proporcional CTS: S/ 1,066',
              '• Total: S/ 15,999',
              '',
              'Genero el asiento de reintegro como concepto REINTEGRO_SUELDO en la',
              'próxima planilla. ¿Apruebas?'],
             font=FONT_BODY, size=11, color=SLATE_900)

    # Botones
    btn_y = chat_y + 3.7
    add_rounded(s, Inches(chat_x + 0.5), Inches(btn_y), Inches(2.5), Inches(0.4), SUCCESS)
    add_text(s, Inches(chat_x + 0.5), Inches(btn_y), Inches(2.5), Inches(0.4),
             '✓ Aprobar reintegro',
             font=FONT_HEAD, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(s, Inches(chat_x + 3.2), Inches(btn_y), Inches(2.2), Inches(0.4), SLATE_300)
    add_text(s, Inches(chat_x + 3.2), Inches(btn_y), Inches(2.2), Inches(0.4),
             'Ver detalle',
             font=FONT_BODY, size=12, color=SLATE_900,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
             'RAG de normativa peruana · 22 entradas (D.S. 003-97-TR, D.Leg. 728/713/650, IR 5ta SUNAT...)',
             font=FONT_BODY, size=12, color=TEAL_LIGHT)
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             'Frase: "El agente actúa, no solo conversa."',
             font=FONT_BODY, size=11, italic=True, color=TEAL_LIGHT)

    page_footer(s, 9, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 10 — Pool de Propinas (gastronomía)
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Pool de Propinas · Gastronomía',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Distribución automática por puntos · sin Excel · sin discusiones',
             font=FONT_BODY, size=14, color=SLATE_500)

    # Tabla puntos
    add_rounded(s, Inches(0.7), Inches(1.85), Inches(6.0), Inches(4.0), SLATE_100, line_color=SLATE_300)
    add_text(s, Inches(0.9), Inches(2.0), Inches(5.5), Inches(0.3),
             'Puntos por cargo (configurable)',
             font=FONT_HEAD, size=14, bold=True, color=SLATE_900)
    cargos = [
        ('Mozo Senior', '3.5'),
        ('Mozo', '3.0'),
        ('Bartender', '2.5'),
        ('Hostess', '1.5'),
        ('Sushi Chef', '2.0'),
        ('Cocinero', '1.0'),
        ('Cocinero Junior', '0.8'),
        ('Ayudante de cocina', '0.5'),
        ('Lavaplatos', '0.5'),
    ]
    for i, (cargo, puntos) in enumerate(cargos):
        y = Inches(2.45 + i * 0.35)
        add_text(s, Inches(0.9), y, Inches(4), Inches(0.3), cargo,
                 font=FONT_BODY, size=11, color=SLATE_700)
        add_text(s, Inches(5.0), y, Inches(1.5), Inches(0.3), puntos + ' pts',
                 font=FONT_MONO, size=11, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

    # Pool ejemplo
    add_rounded(s, Inches(7.0), Inches(1.85), Inches(5.6), Inches(4.0), SLATE_100, line_color=ACCENT)
    add_text(s, Inches(7.2), Inches(2.0), Inches(5), Inches(0.3),
             'Ejemplo: Pool semanal',
             font=FONT_HEAD, size=14, bold=True, color=SLATE_900)
    add_text(s, Inches(7.2), Inches(2.45), Inches(5), Inches(0.4),
             'Monto recolectado: S/ 2,800',
             font=FONT_BODY, size=13, color=SLATE_700)
    add_text(s, Inches(7.2), Inches(2.85), Inches(5), Inches(0.4),
             'Suma puntos del turno: 20.3',
             font=FONT_BODY, size=13, color=SLATE_700)
    add_text(s, Inches(7.2), Inches(3.25), Inches(5), Inches(0.4),
             'Valor por punto: S/ 137.93',
             font=FONT_MONO, size=12, color=ACCENT)

    # Línea separadora
    add_box(s, Inches(7.2), Inches(3.8), Inches(5.2), Inches(0.02), SLATE_300)

    distrib = [
        ('María (Mozo Senior)', '3.5', 'S/ 482.76'),
        ('Carlos (Bartender)', '2.5', 'S/ 344.83'),
        ('Ana (Cocinero)', '1.0', 'S/ 137.93'),
        ('Luis (Ayudante)', '0.5', 'S/ 68.97'),
    ]
    for i, (nombre, pts, monto) in enumerate(distrib):
        y = Inches(3.95 + i * 0.4)
        add_text(s, Inches(7.2), y, Inches(2.8), Inches(0.3), nombre,
                 font=FONT_BODY, size=10, color=SLATE_700)
        add_text(s, Inches(10.0), y, Inches(0.8), Inches(0.3), pts + 'pts',
                 font=FONT_MONO, size=9, color=SLATE_500, align=PP_ALIGN.CENTER)
        add_text(s, Inches(10.8), y, Inches(1.6), Inches(0.3), monto,
                 font=FONT_MONO, size=11, bold=True, color=SUCCESS, align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.3),
             'Investigación: Toast POS, Square, 7shifts, TipHaus, Tronc UK · ' +
             'Modos: POOL_PUNTOS · POOL_PAREJO · INDIVIDUAL · POOL_PUNTOS_HORAS',
             font=FONT_BODY, size=11, color=SLATE_500)
    add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.3),
             'Acta PDF firmable · concepto NO remunerativo (no afecta AFP/CTS/gratis · DS 012-2003-TR)',
             font=FONT_BODY, size=11, color=SLATE_500)
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             'Frase: "Sin Excel, sin discusión entre mozos."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 10, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 11 — Vacaciones + Portal Empleado
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Portal del Empleado',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'RRHH deja de ser el cuello de botella · el trabajador resuelve solo',
             font=FONT_BODY, size=14, color=SLATE_500)

    secciones = [
        ('📅', 'Mi asistencia', 'Tabla mensual con KPIs (T, FA, HE, % asistencia)'),
        ('💰', 'Mis recibos', 'PDFs firmados digitalmente · D.S. 009-2011-TR'),
        ('🏖️', 'Solicitar vacaciones', 'Calendario · días disponibles en tiempo real'),
        ('⏰', 'Banco de horas', 'HE compensatorias acumuladas'),
        ('📋', 'Papeletas y permisos', 'Salida temprana / llegada tarde · aprobación digital'),
        ('🎓', 'Capacitaciones', 'Asistencias · certificados · próximas'),
        ('📂', 'Mis documentos', 'Contrato · constancias · evaluaciones'),
        ('👤', 'Mi perfil', 'Datos personales · contacto · dependientes'),
    ]
    cy = 1.85
    cx = 0.7
    cw = 6.0
    ch = 0.6
    for i, (icon, title, desc) in enumerate(secciones):
        col = i % 2
        row = i // 2
        x = Inches(cx + (cw + 0.3) * col)
        y = Inches(cy + (ch + 0.05) * row)
        add_rounded(s, x, y, Inches(cw), Inches(ch), SLATE_100, line_color=SLATE_300)
        add_text(s, x + Inches(0.15), y, Inches(0.5), Inches(ch), icon,
                 size=18, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.7), y + Inches(0.05), Inches(cw - 0.8), Inches(0.25), title,
                 font=FONT_HEAD, size=12, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.7), y + Inches(0.3), Inches(cw - 0.8), Inches(0.25), desc,
                 font=FONT_BODY, size=10, color=SLATE_500)

    add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
             'Mobile-first · Login con DNI (73000043 / demo123 en demo)',
             font=FONT_BODY, size=11, color=SLATE_500)
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             'Frase: "RRHH deja de ser el cuello de botella."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 11, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 12 — Analytics
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Analytics · People Intelligence',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Predicción, alertas, salud organizacional · "Lo que ningún ERP peruano tiene hoy"',
             font=FONT_BODY, size=14, color=SLATE_500)

    cards = [
        ('🎯', 'Predicción de rotación', 'Score por trabajador · marca riesgo alto/medio/bajo · 3 indicadores claves (pulso, papeletas, tiempo sin promoción)'),
        ('❤️', 'Pulse semanal', 'Encuesta anónima 1 pregunta · score 1-10 · heatmap por local · trend de últimas 12 semanas'),
        ('💵', 'Análisis salarial', 'Bandas por cargo · equidad interna · brecha de género · ajuste sugerido'),
        ('📈', 'Headcount + cost', 'Tendencia última 12 meses · costo planilla · comparativo mes anterior'),
        ('⚠️', 'Riesgo de fuga', 'Top 5 colaboradores con mayor riesgo · razones · plan de acción'),
        ('🎖️', 'Capacitaciones', 'Compliance DIGESA · % de certificados al día · próximos vencimientos'),
    ]
    cx = 0.7
    cy = 1.85
    cw = 4.0
    ch = 1.5
    for i, (icon, title, desc) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = Inches(cx + (cw + 0.2) * col)
        y = Inches(cy + (ch + 0.15) * row)
        add_rounded(s, x, y, Inches(cw), Inches(ch), SLATE_100, line_color=TEAL_LIGHT)
        # icon header
        add_rounded(s, x, y, Inches(cw), Inches(0.45), TEAL)
        add_text(s, x + Inches(0.15), y, Inches(0.4), Inches(0.45), icon,
                 size=18, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.55), y, Inches(cw - 0.6), Inches(0.45), title,
                 font=FONT_HEAD, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(cw - 0.4), Inches(0.85), desc,
                 font=FONT_BODY, size=10, color=SLATE_700)

    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             'Frase: "Esto no lo tiene ningún ERP de RRHH peruano. Hoy."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 12, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 13 — Reclutamiento + CV parser
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Reclutamiento · CV Express',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Pipeline Kanban · CV parser híbrido (local + IA) · cero Excel',
             font=FONT_BODY, size=14, color=SLATE_500)

    # Pipeline kanban
    py = 1.85
    columnas = [
        ('Postulado', '9', RGBColor(0x94, 0xA3, 0xB8)),
        ('Screening', '8', RGBColor(0x3B, 0x82, 0xF6)),
        ('Entrevista', '10', RGBColor(0xF5, 0x9E, 0x0B)),
        ('Prueba', '2', RGBColor(0xA8, 0x55, 0xF7)),
        ('Oferta', '1', RGBColor(0xEC, 0x48, 0x99)),
        ('Contratado', '1', SUCCESS),
    ]
    cw = 2.0
    cx0 = 0.7
    gap = 0.05
    for i, (label, count, color) in enumerate(columnas):
        x = Inches(cx0 + (cw + gap) * i)
        # header
        add_rounded(s, x, Inches(py), Inches(cw), Inches(0.4), color)
        add_text(s, x, Inches(py), Inches(cw), Inches(0.4), label,
                 font=FONT_HEAD, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # count
        add_rounded(s, x, Inches(py + 0.45), Inches(cw), Inches(2.5), SLATE_100)
        add_text(s, x, Inches(py + 0.45), Inches(cw), Inches(1.0), count,
                 font=FONT_MONO, size=36, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, Inches(py + 1.5), Inches(cw), Inches(0.4), 'candidatos',
                 font=FONT_BODY, size=10, color=SLATE_500, align=PP_ALIGN.CENTER)
        if i < 5:
            ax = x + Inches(cw - 0.05)
            add_text(s, ax, Inches(py + 1.6), Inches(0.15), Inches(0.3), '→',
                     size=14, color=SLATE_300, align=PP_ALIGN.CENTER)

    # CV parser feature
    add_rounded(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.8), SLATE_100, line_color=TEAL_LIGHT)
    add_text(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.4),
             '🆕 CV Express · Parser híbrido (Nexotalent-style)',
             font=FONT_HEAD, size=15, bold=True, color=SLATE_900)
    add_text(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.2),
             ['• Fase 1 (LOCAL): regex + spaCy detecta secciones (Datos personales, Experiencia, Educación)',
              '• Fase 2 (IA): solo para enriquecer campos vacíos, no para rehacer todo',
              '• Anti-confusión: NO usa la fecha de nacimiento como inicio de experiencia (bug típico)',
              '• 40 tests passing · Mascarable en CV con sección "Datos personales" + "Fecha nac: 14/03/1990" + experiencia 2018→2026 = 8 años (no 36)'],
             font=FONT_BODY, size=11, color=SLATE_700)

    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             'Frase: "Cero Excel para el reclutador."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 13, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 14 — Workflows + Audit log
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Workflows · Audit · Compliance',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Aprobaciones automáticas · todo lo que pasa, queda',
             font=FONT_BODY, size=14, color=SLATE_500)

    # Workflow ejemplo offboarding
    add_text(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.35),
             'Ejemplo: Workflow Offboarding · 5 etapas en cadena',
             font=FONT_HEAD, size=14, bold=True, color=SLATE_900)

    etapas = ['Encuesta\nsalida', 'Devolución\nactivos', 'Liquidación\npagada', 'Carta\nno adeudo', 'Cierre\nadministrativo']
    ew = 2.2
    ex0 = 0.7
    egap = 0.18
    ey = 2.35
    for i, etapa in enumerate(etapas):
        x = Inches(ex0 + (ew + egap) * i)
        color = TEAL if i < 1 else SLATE_300
        add_rounded(s, x, Inches(ey), Inches(ew), Inches(1.0), color)
        add_text(s, x, Inches(ey), Inches(ew), Inches(1.0), etapa,
                 font=FONT_HEAD, size=12, bold=True, color=WHITE if i < 1 else SLATE_700,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # status badge
        status = 'En curso' if i == 0 else 'Pendiente'
        add_text(s, x, Inches(ey + 1.05), Inches(ew), Inches(0.3), status,
                 font=FONT_BODY, size=10, color=color, align=PP_ALIGN.CENTER)

    # Audit log
    add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.35),
             'Audit Log · 5 modelos críticos con trazabilidad completa',
             font=FONT_HEAD, size=14, bold=True, color=SLATE_900)

    audit_models = [
        ('Personal', 'Altas, bajas, cambios de sueldo, área, jefatura'),
        ('SolicitudVacación', 'Creación, aprobación, rechazo, anulación'),
        ('Préstamo', 'Otorgamiento, cuotas pagadas, refinanciación'),
        ('RegistroPapeleta', 'Creación, aprobación, anulación'),
        ('PeriodoNomina', 'Apertura, cierre, reapertura'),
    ]
    ay = 4.7
    for i, (model, events) in enumerate(audit_models):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + 6.0 * col)
        y = Inches(ay + 0.45 * row)
        add_rounded(s, x, y, Inches(0.7), Inches(0.35), TEAL)
        add_text(s, x, y, Inches(0.7), Inches(0.35), '✓',
                 font=FONT_HEAD, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), y, Inches(2.5), Inches(0.35), model,
                 font=FONT_HEAD, size=12, bold=True, color=SLATE_900,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(3.0), y, Inches(2.8), Inches(0.35), events,
                 font=FONT_BODY, size=10, color=SLATE_500, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.3),
             'Cada cambio queda con usuario · IP · timestamp · diff campo a campo · listo para SUNAFIL',
             font=FONT_BODY, size=11, color=SLATE_500)
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             'Frase: "Si alguien cambia algo, queda registro inmutable."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 14, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 15 — Confiabilidad
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, SLATE_100)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), SUCCESS)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Calidad & Confiabilidad',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'No es promesa · son métricas verificables',
             font=FONT_BODY, size=14, color=SLATE_500)

    # 4 KPIs grandes
    kpi_y = 1.85
    kpi_h = 2.2
    add_kpi_box(s, Inches(0.7),  Inches(kpi_y), Inches(2.85), Inches(kpi_h), '1,816', 'Tests automatizados', SUCCESS)
    add_kpi_box(s, Inches(3.7),  Inches(kpi_y), Inches(2.85), Inches(kpi_h), '94%', 'Cobertura HE engine', TEAL)
    add_kpi_box(s, Inches(6.7),  Inches(kpi_y), Inches(2.85), Inches(kpi_h), '91%', 'Cobertura nóminas', TEAL)
    add_kpi_box(s, Inches(9.7),  Inches(kpi_y), Inches(2.85), Inches(kpi_h), '12 min', 'Cierre planilla', ACCENT)

    # Lista de garantías
    garantias = [
        ('🔒', 'Backup automático diario', 'pg_dump 03:30 AM · retención 30 días · S3/Backblaze B2'),
        ('🚨', 'Sentry para errores producción', 'Tags por tenant · breadcrumbs por usuario · alertas inmediatas'),
        ('📜', 'Audit log completo', '5 modelos críticos · usuario · IP · diff campo a campo'),
        ('🏛️', 'Normativa peruana centralizada', 'D.S. 003-97-TR · D.Leg. 728/713/650/1057 · Ley 26790/27735/30367'),
    ]
    gy = 4.3
    for i, (icon, title, desc) in enumerate(garantias):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + 6.0 * col)
        y = Inches(gy + 0.85 * row)
        add_rounded(s, x, y, Inches(5.8), Inches(0.75), WHITE, line_color=SLATE_300)
        add_text(s, x + Inches(0.15), y, Inches(0.7), Inches(0.75), icon,
                 size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.9), y + Inches(0.08), Inches(4.8), Inches(0.35), title,
                 font=FONT_HEAD, size=13, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.9), y + Inches(0.4), Inches(4.8), Inches(0.35), desc,
                 font=FONT_BODY, size=10, color=SLATE_500)

    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             'Frase: "Lo detectamos antes que ustedes."',
             font=FONT_BODY, size=11, italic=True, color=TEAL)

    page_footer(s, 15, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 16 — Comparativa
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             '¿Por qué Harmoni vs alternativas?',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'Comparativa cabeza a cabeza',
             font=FONT_BODY, size=14, color=SLATE_500)

    # Tabla 2 columnas
    # Headers
    add_rounded(s, Inches(0.7),  Inches(1.85), Inches(6.2), Inches(0.45), SLATE_500)
    add_text(s, Inches(0.7), Inches(1.85), Inches(6.2), Inches(0.45),
             'Competidor típico (SAP / Bambu / Workday)',
             font=FONT_HEAD, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(s, Inches(6.95), Inches(1.85), Inches(5.65), Inches(0.45), TEAL)
    add_text(s, Inches(6.95), Inches(1.85), Inches(5.65), Inches(0.45),
             'Harmoni',
             font=FONT_HEAD, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    comparativa = [
        ('Sistema importado, español neutro', 'Hecho en Perú, normativa peruana adentro'),
        ('Cierre planilla 2-4 horas con Excel paralelo', '12 minutos en wizard guiado'),
        ('Asistente que conversa', 'Agente IA que ACTÚA con tool-use real'),
        ('Sin portal del empleado', 'Portal con auto-servicio mobile-first'),
        ('Boleta PDF improvisada', 'DS 009-2011-TR · firma electrónica · acuse'),
        ('Sin audit log', 'AuditEntry en 5 modelos con diff campo a campo'),
        ('Sin backup automático', 'pg_dump diario · retención 30d · S3/B2'),
        ('Sin tests', '1,816 tests · 94% cobertura HE'),
        ('Diseño anticuado SAP-style', 'Cockpit V2 · Linear / Notion-style'),
    ]
    ry = 2.35
    rh = 0.45
    for i, (comp, harmoni) in enumerate(comparativa):
        y = Inches(ry + rh * i)
        bg_col = SLATE_100 if i % 2 == 0 else WHITE
        add_box(s, Inches(0.7),  y, Inches(6.2), Inches(rh), bg_col)
        add_box(s, Inches(6.95), y, Inches(5.65), Inches(rh), bg_col)
        add_text(s, Inches(0.9), y, Inches(5.9), Inches(rh), '✗ ' + comp,
                 font=FONT_BODY, size=10, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.15), y, Inches(5.35), Inches(rh), '✓ ' + harmoni,
                 font=FONT_HEAD, size=10, bold=True, color=SUCCESS, anchor=MSO_ANCHOR.MIDDLE)

    page_footer(s, 16, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 17 — Roadmap
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)

    add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.5),
             'Roadmap',
             font=FONT_HEAD, size=36, bold=True, color=SLATE_900)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
             'En producción hoy · Próximas 6 semanas · Q3 2026',
             font=FONT_BODY, size=14, color=SLATE_500)

    # 3 columnas
    cols = [
        ('✅ Hoy en producción', SUCCESS, [
            'Personal V2 (drawer + bulk + inline)',
            'Asistencia matrix + sync biométrico',
            'Nóminas wizard 11 pasos',
            'Liquidaciones unificadas (Sprint 1-3)',
            'Workflow offboarding 5 etapas',
            'Pool propinas gastronomía',
            'Agente IA con RAG normativa',
            'Portal del empleado',
            'Analytics + Predictive',
            'Audit log v2',
            'Backup automatizado',
        ]),
        ('🟡 Próximas 6 semanas', WARNING, [
            'Sprint 4 Liquidaciones · E2E tests + capacitación cliente',
            'Pool propinas: anomaly detection',
            'Auto-import propinas desde POS (R&G, Sysrest)',
            'Pay-card estilo HausMoney/Kickfin',
            'App móvil portal del empleado',
            'Workflow visual editor (drag-drop)',
        ]),
        ('🔵 Q3 2026', TEAL, [
            'Multi-empresa selector global',
            'IA generadora de reportes ejecutivos',
            'Integración SUNAT directa AFP Net',
            'Roles finos por módulo',
            'Limpieza datos legacy CSRT',
            'Dashboard cliente personalizable',
            'Marketplace plantillas',
        ]),
    ]
    cw = 4.0
    cx0 = 0.7
    gap = 0.15
    cy = 1.85
    for i, (titulo, color, items) in enumerate(cols):
        x = Inches(cx0 + (cw + gap) * i)
        add_rounded(s, x, Inches(cy), Inches(cw), Inches(5.0), SLATE_100, line_color=color)
        add_rounded(s, x, Inches(cy), Inches(cw), Inches(0.5), color)
        add_text(s, x, Inches(cy), Inches(cw), Inches(0.5), titulo,
                 font=FONT_HEAD, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # items
        iy = cy + 0.65
        for item in items:
            add_text(s, x + Inches(0.2), Inches(iy), Inches(cw - 0.4), Inches(0.3),
                     '• ' + item,
                     font=FONT_BODY, size=10, color=SLATE_700)
            iy += 0.35

    page_footer(s, 17, TOTAL)

    # ─────────────────────────────────────────────────────────────────
    # Slide 18 — Cierre
    # ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    set_bg(s, SLATE_900)
    add_box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL_LIGHT)

    add_text(s, Inches(1), Inches(1.2), Inches(11.3), Inches(0.8),
             '¿Qué tres procesos',
             font=FONT_HEAD, size=48, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
             'de su RRHH están',
             font=FONT_HEAD, size=48, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(2.8), Inches(11.3), Inches(0.8),
             'consumiendo más tiempo este mes?',
             font=FONT_HEAD, size=48, bold=True, color=TEAL_LIGHT,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
             'Empezamos por ahí.',
             font=FONT_BODY, size=24, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    add_rounded(s, Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.6), TEAL)
    add_text(s, Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.6),
             'demo.harmoni.pe',
             font=FONT_MONO, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
             'Edwin Lopez · eflopezt@gmail.com · Harmoni · Mayo 2026',
             font=FONT_BODY, size=12, color=SLATE_300,
             align=PP_ALIGN.CENTER)

    # ─────────────────────────────────────────────────────────────────
    # Guardar
    # ─────────────────────────────────────────────────────────────────
    out_dir = Path('presentacion')
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / 'Harmoni_Demo_Mayo_2026.pptx'
    prs.save(str(out_path))
    print(f'\n[OK] PPT generado: {out_path}')
    print(f'  Slides: {len(prs.slides)}')
    print(f'  Tamano: {os.path.getsize(out_path) / 1024:.1f} KB')


if __name__ == '__main__':
    main()
