"""
Presentacion de Ventas — Harmoni ERP (PPTX)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

IMGS = r"D:\Harmoni\docs\manual_capturas"
TEAL = RGBColor(0x0F, 0x76, 0x6E)
DARK = RGBColor(0x0A, 0x2F, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
CYAN = RGBColor(0x5E, 0xEA, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_multiline(slide, left, top, width, height, lines, size=14, color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(6)
    return tf

def add_img(slide, path, left, top, width=None, height=None):
    full = os.path.join(IMGS, path)
    if os.path.exists(full):
        if width:
            slide.shapes.add_picture(full, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(full, Inches(left), Inches(top), height=Inches(height))

# ================================================================
# SLIDE 1 - PORTADA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_text(slide, 1, 1.5, 11, 1, 'H  HARMONI', size=60, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3, 11, 1, 'El ERP de RRHH que tu empresa necesita', size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.2, 11, 0.8, '22 Modulos  |  100% Normativa Peruana  |  IA Integrada', size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.6, 'www.harmoni.pe  |  +51 977 538 028', size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 2 - PROBLEMA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, 0.8, 0.3, 12, 0.8, 'El problema', size=36, color=DARK, bold=True)

problemas = [
    'Planilla en Excel con riesgo de errores y multas SUNAFIL',
    'Contratos vencidos sin alertas — contingencias legales',
    'Horas extra sin control — perdida de dinero',
    'Vacaciones calculadas "a ojo" — incumplimiento DL 713',
    'Exportaciones a SUNAT/AFP manuales — horas desperdiciadas',
    'Sin visibilidad de tu fuerza laboral en tiempo real',
]
add_multiline(slide, 0.8, 1.3, 6, 5, problemas, size=16, color=RGBColor(0xC0, 0x00, 0x00))

add_text(slide, 7.5, 0.3, 5, 0.8, 'La solucion: Harmoni', size=28, color=TEAL, bold=True)
soluciones = [
    'Motor de nominas automatico (IR 5ta, AFP, CTS)',
    'Contratos DL 728 con alertas y PDF/DOCX',
    'Asistencia biometrica con reglas inteligentes',
    'Vacaciones DL 713 al dia con saldos y vencidas',
    'PLAME, T-Registro, AFP Net en un click',
    'Dashboard en tiempo real con 200+ KPIs',
]
add_multiline(slide, 7.5, 1.3, 5.5, 5, soluciones, size=16, color=TEAL)

# ================================================================
# SLIDE 3 - DASHBOARD
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF9, 0xFA, 0xFB))
add_text(slide, 0.5, 0.2, 12, 0.7, 'Dashboard Principal', size=30, color=DARK, bold=True)
add_img(slide, '01_dashboard/01_main.png', 0.5, 1, width=12.3)

# ================================================================
# SLIDE 4 - PERSONAL
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, 0.5, 0.2, 6, 0.7, 'Gestion de Personal', size=30, color=DARK, bold=True)
add_multiline(slide, 0.5, 1, 5.5, 5, [
    'Ficha completa de 50+ campos por empleado',
    'Organigrama interactivo con colores por nivel',
    'Contratos DL 728: PDF, DOCX, prorrogas, adendas',
    'Areas, Sub-Areas, Cargos normalizados',
    'Roster matricial de asignaciones',
    'Reportes: headcount, rotacion, contratos',
    'Importacion masiva desde Excel',
], size=14, color=GRAY)
add_img(slide, '02_personal/06_organigrama.png', 6.5, 0.5, width=6.3)

# ================================================================
# SLIDE 5 - ASISTENCIA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF9, 0xFA, 0xFB))
add_text(slide, 0.5, 0.2, 12, 0.7, 'Control de Asistencia', size=30, color=DARK, bold=True)
add_img(slide, '03_asistencia/03_calendario_marzo.png', 0.3, 1, width=12.7)

# ================================================================
# SLIDE 6 - REPORTES STAFF/RCO
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, 0.5, 0.2, 5, 0.7, 'Reporte Staff', size=24, color=DARK, bold=True)
add_img(slide, '03_asistencia/07_reporte_staff.png', 0.2, 0.9, width=6.3)
add_text(slide, 7, 0.2, 5, 0.7, 'Reporte RCO', size=24, color=DARK, bold=True)
add_img(slide, '03_asistencia/08_reporte_rco.png', 6.8, 0.9, width=6.3)

# ================================================================
# SLIDE 7 - CONTRATOS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF9, 0xFA, 0xFB))
add_text(slide, 0.5, 0.2, 12, 0.7, 'Contratos Laborales — DL 728', size=30, color=DARK, bold=True)
add_img(slide, '02_personal/09_contratos_panel.png', 0.3, 1, width=12.7)

# ================================================================
# SLIDE 8 - NOMINAS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, 0.5, 0.2, 12, 0.7, 'Nominas y Planillas', size=30, color=DARK, bold=True)
add_text(slide, 0.5, 0.8, 5, 0.5, 'Calculo automatico de todos los conceptos', size=14, color=GRAY)
add_multiline(slide, 0.5, 1.5, 5.5, 5, [
    'Sueldo basico + asignacion familiar',
    'Horas extra: 25%, 35%, 100%',
    'Gratificacion + bonificacion 9%',
    'CTS semestral',
    'AFP (Habitat, Integra, Prima, Profuturo)',
    'ONP',
    'IR 5ta categoria (escala progresiva)',
    'Descuentos: prestamos, adelantos, judiciales',
    'Boletas PDF individuales y masivas',
    'Liquidacion automatica por cese',
], size=13, color=GRAY)
add_img(slide, '04_nominas/01_periodos.png', 6.5, 0.8, width=6.3)

# ================================================================
# SLIDE 9 - 22 MODULOS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, 0.5, 0.2, 12, 0.8, '22 Modulos Integrados', size=36, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

mods_left = [
    'Personal & Ficha Empleado',
    'Organigrama Interactivo',
    'Contratos DL 728',
    'Asistencia Matricial',
    'Papeletas & Permisos',
    'Banco de Horas Extra',
    'Nominas & Planillas',
    'Boletas PDF',
    'Liquidaciones',
    'Vacaciones DL 713',
    'IR 5ta Categoria',
]
mods_right = [
    'Reclutamiento Kanban',
    'Onboarding & Offboarding',
    'Evaluaciones 360 / OKR',
    'Capacitaciones & LMS',
    'Encuestas de Clima',
    'Disciplinaria',
    'Prestamos',
    'Viaticos',
    'Comunicaciones & WhatsApp',
    'Documentos & Firma Digital',
    'Workflows & Aprobaciones',
]

add_multiline(slide, 0.8, 1.2, 5.5, 5.5, mods_left, size=14, color=WHITE)
add_multiline(slide, 7, 1.2, 5.5, 5.5, mods_right, size=14, color=WHITE)

# ================================================================
# SLIDE 10 - IA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, 0.5, 0.3, 12, 0.8, 'Inteligencia Artificial Integrada', size=36, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0.5, 1.2, 12, 0.6, 'La unica plataforma de RRHH en Peru con IA nativa', size=18, color=WHITE, align=PP_ALIGN.CENTER)

ia = [
    'Asistente Virtual 24/7 — consultas de personal, normativa, asistencia',
    'Analisis de Contratos — riesgos legales, cumplimiento DL 728',
    'Generacion Automatica — contratos, adendas, prorrogas desde plantillas',
    'Base de Conocimiento RAG — 59 articulos de normativa peruana',
    'Reglas Inteligentes — configura asistencia conversando con la IA',
    'Dashboard Predictivo — rotacion, riesgo de fuga, clima laboral',
    'Multi-proveedor — Gemini 2.5 + DeepSeek V3 + GPT-4o',
]
add_multiline(slide, 1.5, 2.2, 10, 4.5, ia, size=16, color=RGBColor(0xC8, 0xFA, 0xF2))

# ================================================================
# SLIDE 11 - INTEGRACIONES
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, 0.5, 0.2, 12, 0.7, 'Integraciones Peruanas — Un Click', size=30, color=DARK, bold=True, align=PP_ALIGN.CENTER)

integ = [
    'SUNAT: PLAME mensual + T-Registro altas/bajas + PDT 601',
    'AFP: Habitat, Integra, Prima, Profuturo (formato AFP Net)',
    'Bancos: Telecredito BCP, BBVA, Interbank, Scotiabank',
    'CTS: Depositos semestrales en formato de cada banco',
    'EsSalud: Aportes patronales mensuales',
    'Contabilidad: CONCAR, SIGO, SAP, SIRE',
    'SCTR: Seguro Complementario de Trabajo de Riesgo',
]
add_multiline(slide, 1.5, 1.2, 10, 4, integ, size=16, color=GRAY)
add_img(slide, '13_integraciones/01_panel.png', 1, 4.5, width=11)

# ================================================================
# SLIDE 12 - PORTAL
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF9, 0xFA, 0xFB))
add_text(slide, 0.5, 0.2, 12, 0.7, 'Portal del Trabajador — Autoservicio', size=30, color=DARK, bold=True)
add_img(slide, '15_portal/01_portal_main.png', 0.3, 1, width=12.7)

# ================================================================
# SLIDE 13 - PRECIOS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, 0.5, 0.3, 12, 0.8, 'Planes y Precios', size=36, color=DARK, bold=True, align=PP_ALIGN.CENTER)

# Starter
add_text(slide, 1, 1.5, 3.5, 0.6, 'STARTER', size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 2.2, 3.5, 0.6, 'S/ 149/mes', size=28, color=DARK, bold=True, align=PP_ALIGN.CENTER)
add_multiline(slide, 1, 3, 3.5, 3.5, [
    'Hasta 50 colaboradores',
    'Personal + Asistencia',
    'Nominas basicas',
    'Portal del Trabajador',
    'Soporte email',
], size=13, color=GRAY)

# Profesional (featured)
add_text(slide, 5, 1.2, 3.5, 0.4, 'MAS POPULAR', size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 5, 1.5, 3.5, 0.6, 'PROFESIONAL', size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 5, 2.2, 3.5, 0.6, 'S/ 349/mes', size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_multiline(slide, 5, 3, 3.5, 3.5, [
    'Hasta 200 colaboradores',
    '22 modulos completos',
    'Inteligencia Artificial',
    'PLAME, AFP, CTS bancos',
    'Soporte WhatsApp prioritario',
    'Implementacion asistida',
], size=13, color=RGBColor(0xE8, 0xFA, 0xF7))

# Paint the Profesional column
txBox = slide.shapes.add_textbox(Inches(4.8), Inches(1.1), Inches(3.9), Inches(5.5))
tf = txBox.text_frame
# Can't easily paint bg in pptx textbox, skip

# Enterprise
add_text(slide, 9, 1.5, 3.5, 0.6, 'ENTERPRISE', size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 9, 2.2, 3.5, 0.6, 'A medida', size=28, color=DARK, bold=True, align=PP_ALIGN.CENTER)
add_multiline(slide, 9, 3, 3.5, 3.5, [
    'Colaboradores ilimitados',
    'Multi-empresa',
    'API REST completa',
    'Personalizaciones',
    'Servidor dedicado',
    'SLA 24/7',
], size=13, color=GRAY)

# ================================================================
# SLIDE 14 - CTA
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, 1, 1.5, 11, 1, 'Empieza hoy', size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3, 11, 0.8, 'Solicita tu demo gratuita', size=28, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11, 0.6, 'www.harmoni.pe', size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.3, 11, 0.6, '+51 977 538 028  |  ventas@harmoni.pe', size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6.3, 11, 0.5, 'Harmoni ERP  |  Lima, Peru  |  2026', size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ================================================================
# GUARDAR
# ================================================================
out = r"D:\Harmoni\docs\PRESENTACION VENTAS - HARMONI ERP.pptx"
prs.save(out)
print(f"PPTX: {out}")
print(f"Total slides: {len(prs.slides)}")
print("OK")
